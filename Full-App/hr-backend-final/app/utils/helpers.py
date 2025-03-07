import io
import os
import json
import docx
import httpx
import openai
import PyPDF2
import numpy as np
from typing import Dict
from docx import Document
from pptx import Presentation
from app.core.config import get_settings
from pdfminer.high_level import extract_text
from langchain.chat_models import ChatOpenAI


model = ChatOpenAI(model="gpt-4o-mini")


document_store: Dict[str, dict] = {}

settings = get_settings()


def parse_request(data: str):
    try:
        return json.loads(data)
    except json.JSONDecodeError:
        return {"error": "Invalid JSON format"}


def get_embedding(text: str):
    """Generate text embeddings using OpenAI API."""
    response = openai.Embedding.create(
        input=text,
        model="text-embedding-ada-002",
        api_key=settings.OPENAI_API_KEY
    )
    return response["data"][0]["embedding"]


def index_all_documents():
    docs_folder = "../docs"
    if not os.path.exists(docs_folder):
        os.makedirs(docs_folder)

    for file_name in os.listdir(docs_folder):
        file_path = os.path.join(docs_folder, file_name)
        if file_name not in document_store:
            process_and_index_file(file_path, file_name)



def process_and_index_file(file_path: str, doc_id: str):
    try:
        file_extension = os.path.splitext(file_path)[-1].lower()
        if file_extension == ".pdf":
            text = extract_text(file_path)
        elif file_extension == ".docx":
            text = extract_text_from_docx(file_path)
        elif file_extension == ".pptx":
            text = extract_text_from_pptx(file_path)
        elif file_extension == ".txt":
            text = extract_text_from_txt(file_path)
        else:
            raise ValueError("Unsupported file format")

        if text:
            embedding = get_embedding(text)
            document_store[doc_id] = {"text": text, "embedding": embedding}
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        raise


def extract_text_from_docx(docx_path: str) -> str:
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                text_content.append(shape.text_frame.text)
    return "\n".join(text_content)


def extract_text_from_txt(txt_path: str) -> str:
    with open(txt_path, "r", encoding="utf-8") as file:
        return file.read()


async def analyze_text(text: str, prompt: str,model:str):
    try:
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {settings.OPENAI_API_KEY}"
        }

        payload = {
            "model":model.lower(),
            "messages": [
                {"role": "system", "content": "You are a helpful AI assistant that analyzes documents."},
                {"role": "user", "content": f"Answer this question {prompt} based on this: {text}"}
            ],
            "stream": True
        }

        async with httpx.AsyncClient() as client:
            async with client.stream(
                "POST",
                "https://api.openai.com/v1/chat/completions",
                headers=headers,
                json=payload,
                timeout=30.0
            ) as response:
                if response.status_code != 200:
                    yield f"Error: API request failed with status {response.status_code}"
                    return

                async for line in response.aiter_lines():
                    if line.startswith('data: '):
                        line = line[6:]
                        if line.strip() == '[DONE]':
                            continue
                        try:
                            data = json.loads(line)
                            if data.get('choices') and data['choices'][0].get('delta', {}).get('content'):
                                yield data['choices'][0]['delta']['content']
                        except json.JSONDecodeError:
                            continue

    except Exception as e:
        yield f"Error during text analysis: {str(e)}"


def extract_text_from_pdf(file: io.BytesIO) -> str:
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text


def extract_text_from_docx(file: io.BytesIO) -> str:
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def extract_text_from_txt(file: io.BytesIO) -> str:
    return file.read().decode("utf-8")
